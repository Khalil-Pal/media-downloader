"""
services/converter.py - File conversion logic.

This module is the only place that shells out to ffmpeg/LibreOffice.
"""
from __future__ import annotations

import asyncio
import logging
import shutil
import uuid
from dataclasses import dataclass
from pathlib import Path

from config.settings import settings

logger = logging.getLogger(__name__)

CONVERSION_TIMEOUT_SECONDS = 120
SLOW_CONVERSION_TIMEOUT_SECONDS = 180
IMAGE_FORMATS = {"png", "jpg", "jpeg", "webp"}
OFFICE_TO_PDF = {"docx", "xlsx", "pptx"}
MARKDOWN_FORMATS = {"md", "markdown"}
PDF_FORMATS = {"pdf"}
PDF_TO_OFFICE = {"docx", "pptx"}
PDF_TO_TEXT_FORMATS = {"md", "xlsx"}
DOCX_TO_PPTX = ("docx", "pptx")
AUDIO_VIDEO_EXTENSIONS = {
    "3gp", "aac", "aiff", "ape", "avi", "flac", "flv", "m4a", "m4v",
    "mkv", "mov", "mp3", "mp4", "mpeg", "mpg", "oga", "ogg", "opus",
    "ts", "wav", "webm", "wma", "wmv",
}

_conversion_semaphore = asyncio.Semaphore(settings.max_concurrent_conversions)


@dataclass(frozen=True)
class ConversionOption:
    target_format: str
    label: str


class ConversionError(Exception):
    """Known conversion failure mapped to a translated user-facing message."""

    def __init__(self, code: str) -> None:
        super().__init__(code)
        self.code = code


def _extension(filename: str | None) -> str:
    if not filename:
        return ""
    suffix = Path(filename).suffix.lower().lstrip(".")
    return "jpg" if suffix == "jpeg" else suffix


def supported_targets(filename: str | None, mime_type: str | None = None) -> list[ConversionOption]:
    """Return supported target formats for a file name / MIME type."""
    ext = _extension(filename)
    mime = (mime_type or "").lower()

    if ext in PDF_FORMATS or mime == "application/pdf":
        return [
            ConversionOption("docx", "-> Word"),
            ConversionOption("pptx", "-> PowerPoint"),
            ConversionOption("md", "-> Markdown"),
            ConversionOption("xlsx", "-> Excel"),
        ]

    if ext in IMAGE_FORMATS:
        source = "jpg" if ext == "jpeg" else ext
        return [
            ConversionOption(target, f"-> {target.upper()}")
            for target in ("png", "jpg", "webp")
            if target != source
        ]

    if ext in OFFICE_TO_PDF or ext in MARKDOWN_FORMATS:
        options = [ConversionOption("pdf", "-> PDF")]
        if ext == "docx":
            options.append(ConversionOption("pptx", "-> PowerPoint"))
        return options

    if (mime.startswith("audio/") or mime.startswith("video/") or ext in AUDIO_VIDEO_EXTENSIONS) and ext != "mp3":
        return [ConversionOption("mp3", "-> MP3")]

    return []


def create_conversion_session() -> Path:
    session_dir = settings.download_path / str(uuid.uuid4())
    session_dir.mkdir(parents=True, exist_ok=True)
    return session_dir


def cleanup_conversion_session(path: Path | None) -> None:
    if not path:
        return

    try:
        download_root = settings.download_path.resolve()
        target = path.resolve()
        if target == download_root or download_root not in target.parents:
            logger.warning("Refusing conversion cleanup outside download directory: %s", path)
            return
        if target.exists():
            shutil.rmtree(target)
    except OSError as exc:
        logger.warning("Conversion cleanup failed: %s", exc)


def friendly_conversion_error(exc: Exception) -> str:
    if isinstance(exc, ConversionError):
        return {
            "too_large": "conversion_error_too_large",
            "timeout": "conversion_error_timeout",
            "unsupported": "conversion_error_unsupported",
            "missing_tool": "conversion_error_missing_tool",
            "tool_failed": "conversion_error_failed",
            "empty_output": "conversion_error_failed",
            "scanned_pdf": "conversion_error_pdf_scanned",
            "no_tables": "conversion_error_pdf_no_tables",
            "pdf_text_failed": "conversion_error_pdf_text_failed",
        }.get(exc.code, "conversion_error_failed")
    return "conversion_error_failed"


def tier2_caveat_key(
    filename: str | None,
    target_format: str,
    mime_type: str | None = None,
) -> str | None:
    source_ext = _extension(filename)
    source_is_pdf = source_ext == "pdf" or (mime_type or "").lower() == "application/pdf"
    target = target_format.lower().lstrip(".")
    if (source_ext, target) == DOCX_TO_PPTX:
        return "conversion_tier2_caveat_docx_pptx"
    if source_is_pdf and target in (PDF_TO_OFFICE | PDF_TO_TEXT_FORMATS):
        return "conversion_tier2_caveat"
    return None


async def convert_file(input_path: Path, target_format: str, mime_type: str | None = None) -> Path:
    """Convert one local file to one supported target format."""
    target = target_format.lower().lstrip(".")
    if target == "jpeg":
        target = "jpg"

    async with _conversion_semaphore:
        if input_path.stat().st_size > settings.max_convert_file_size_bytes:
            raise ConversionError("too_large")

        targets = {option.target_format for option in supported_targets(input_path.name, mime_type)}
        if target not in targets:
            raise ConversionError("unsupported")

        source_ext = _extension(input_path.name)
        source_is_pdf = source_ext == "pdf" or (mime_type or "").lower() == "application/pdf"
        output_path = input_path.with_name(f"{input_path.stem}.{target}")

        if target == "mp3":
            await _convert_to_mp3(input_path, output_path)
        elif source_ext in IMAGE_FORMATS and target in IMAGE_FORMATS:
            await asyncio.to_thread(_convert_image, input_path, output_path, target)
        elif source_ext in OFFICE_TO_PDF and target == "pdf":
            output_path = await _convert_with_libreoffice(input_path, "pdf")
        elif source_ext in MARKDOWN_FORMATS and target == "pdf":
            await asyncio.to_thread(_convert_markdown_to_pdf, input_path, output_path)
        elif source_is_pdf and target in PDF_TO_OFFICE:
            _ensure_pdf_has_text(input_path)
            # LibreOffice PDF import has low fidelity for complex layouts, tables,
            # and scanned/image-only documents.
            output_path = await _convert_with_libreoffice(
                input_path,
                target,
                timeout=SLOW_CONVERSION_TIMEOUT_SECONDS,
            )
        elif source_is_pdf and target == "md":
            _ensure_pdf_has_text(input_path)
            await _run_threaded_conversion(
                _convert_pdf_to_markdown,
                input_path,
                output_path,
                timeout=SLOW_CONVERSION_TIMEOUT_SECONDS,
            )
        elif source_is_pdf and target == "xlsx":
            _ensure_pdf_has_text(input_path)
            await _run_threaded_conversion(
                _convert_pdf_to_xlsx,
                input_path,
                output_path,
                timeout=SLOW_CONVERSION_TIMEOUT_SECONDS,
            )
        elif (source_ext, target) == DOCX_TO_PPTX:
            await _run_threaded_conversion(
                _convert_docx_to_pptx,
                input_path,
                output_path,
                timeout=SLOW_CONVERSION_TIMEOUT_SECONDS,
            )
        else:
            raise ConversionError("unsupported")

        if not output_path.exists() or output_path.stat().st_size == 0:
            raise ConversionError("empty_output")
        return output_path


async def _run_process(
    command: list[str],
    tool_name: str,
    timeout: int = CONVERSION_TIMEOUT_SECONDS,
) -> None:
    try:
        process = await asyncio.create_subprocess_exec(
            *command,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
    except FileNotFoundError as exc:
        logger.warning("%s is not installed or not on PATH.", tool_name)
        raise ConversionError("missing_tool") from exc

    try:
        stdout, stderr = await asyncio.wait_for(
            process.communicate(),
            timeout=timeout,
        )
    except asyncio.TimeoutError as exc:
        process.kill()
        await process.communicate()
        logger.warning("%s conversion timed out.", tool_name)
        raise ConversionError("timeout") from exc

    if process.returncode != 0:
        details = (stderr or stdout or b"").decode("utf-8", errors="replace").strip()
        logger.warning("%s conversion failed: %s", tool_name, details[-1000:])
        raise ConversionError("tool_failed")


async def _run_threaded_conversion(function, *args, timeout: int) -> None:
    try:
        await asyncio.wait_for(asyncio.to_thread(function, *args), timeout=timeout)
    except asyncio.TimeoutError as exc:
        raise ConversionError("timeout") from exc


async def _convert_to_mp3(input_path: Path, output_path: Path) -> None:
    output_path.unlink(missing_ok=True)
    await _run_process(
        [
            "ffmpeg", "-y", "-hide_banner", "-loglevel", "error",
            "-i", str(input_path),
            "-vn", "-codec:a", "libmp3lame", "-b:a", "192k",
            str(output_path),
        ],
        "ffmpeg",
    )


def _convert_image(input_path: Path, output_path: Path, target_format: str) -> None:
    try:
        from PIL import Image
    except ModuleNotFoundError as exc:
        raise ConversionError("missing_tool") from exc

    output_path.unlink(missing_ok=True)
    with Image.open(input_path) as image:
        if target_format == "jpg":
            if image.mode in {"RGBA", "LA"} or "transparency" in image.info:
                rgba = image.convert("RGBA")
                background = Image.new("RGBA", rgba.size, (255, 255, 255, 255))
                background.alpha_composite(rgba)
                converted = background.convert("RGB")
            else:
                converted = image.convert("RGB")
            converted.save(output_path, format="JPEG", quality=95, optimize=True)
            return

        save_format = "PNG" if target_format == "png" else "WEBP"
        image.save(output_path, format=save_format)


async def _convert_with_libreoffice(
    input_path: Path,
    target_format: str,
    timeout: int = CONVERSION_TIMEOUT_SECONDS,
) -> Path:
    output_path = input_path.with_suffix("." + target_format)
    output_path.unlink(missing_ok=True)
    profile_dir = input_path.parent / ("lo_profile_" + uuid.uuid4().hex)
    profile_dir.mkdir(parents=True, exist_ok=True)
    try:
        await _run_process(
            [
                "soffice",
                f"-env:UserInstallation={profile_dir.resolve().as_uri()}",
                "--headless",
                "--convert-to", target_format,
                "--outdir", str(input_path.parent),
                str(input_path),
            ],
            "LibreOffice",
            timeout=timeout,
        )
    finally:
        shutil.rmtree(profile_dir, ignore_errors=True)
    return output_path


def _pdf_text_stats(input_path: Path) -> tuple[int, int]:
    try:
        import fitz
    except ModuleNotFoundError as exc:
        raise ConversionError("missing_tool") from exc

    try:
        with fitz.open(input_path) as document:
            page_count = document.page_count
            text_length = sum(
                len(page.get_text("text").strip())
                for page in document
            )
    except Exception as exc:
        logger.warning("PDF text extraction failed for %s: %s", input_path, exc)
        raise ConversionError("pdf_text_failed") from exc

    return page_count, text_length


def _ensure_pdf_has_text(input_path: Path) -> None:
    page_count, text_length = _pdf_text_stats(input_path)
    if page_count and text_length < max(20, page_count * 20):
        raise ConversionError("scanned_pdf")


def _convert_markdown_to_pdf(input_path: Path, output_path: Path) -> None:
    try:
        import markdown
        from weasyprint import HTML
    except ModuleNotFoundError as exc:
        raise ConversionError("missing_tool") from exc

    output_path.unlink(missing_ok=True)
    source = input_path.read_text(encoding="utf-8", errors="replace")
    body = markdown.markdown(source, extensions=["extra", "sane_lists", "tables"])
    html = (
        "<!doctype html><html><head><meta charset='utf-8'>"
        "<style>body{font-family:Arial,sans-serif;line-height:1.5;margin:32px;}"
        "pre,code{font-family:Consolas,monospace;} img{max-width:100%;}</style>"
        "</head><body>" + body + "</body></html>"
    )
    HTML(string=html, base_url=str(input_path.parent)).write_pdf(str(output_path))


def _convert_pdf_to_markdown(input_path: Path, output_path: Path) -> None:
    output_path.unlink(missing_ok=True)
    try:
        import pymupdf4llm

        markdown_text = pymupdf4llm.to_markdown(str(input_path))
    except Exception as exc:
        logger.warning("pymupdf4llm failed for %s; falling back to plain text: %s", input_path, exc)
        markdown_text = _extract_pdf_plain_markdown(input_path)

    if not markdown_text.strip():
        raise ConversionError("scanned_pdf")
    output_path.write_text(markdown_text, encoding="utf-8")


def _extract_pdf_plain_markdown(input_path: Path) -> str:
    try:
        import fitz
    except ModuleNotFoundError as exc:
        raise ConversionError("missing_tool") from exc

    try:
        chunks: list[str] = []
        with fitz.open(input_path) as document:
            for index, page in enumerate(document, start=1):
                text = page.get_text("text").strip()
                if text:
                    chunks.append(f"## Page {index}\n\n{text}")
    except Exception as exc:
        logger.warning("Plain PDF text fallback failed for %s: %s", input_path, exc)
        raise ConversionError("pdf_text_failed") from exc
    return "\n\n".join(chunks)


def _convert_pdf_to_xlsx(input_path: Path, output_path: Path) -> None:
    try:
        import camelot
        from openpyxl import Workbook
    except ModuleNotFoundError as exc:
        raise ConversionError("missing_tool") from exc

    output_path.unlink(missing_ok=True)
    try:
        tables = camelot.read_pdf(str(input_path), pages="all", flavor="stream")
    except Exception as exc:
        logger.warning("Camelot table extraction failed for %s: %s", input_path, exc)
        raise ConversionError("tool_failed") from exc

    usable_tables = [
        table for table in tables
        if table.df.shape[0] >= 2 and table.df.shape[1] >= 2
    ]
    if not usable_tables:
        raise ConversionError("no_tables")

    workbook = Workbook()
    workbook.remove(workbook.active)
    for index, table in enumerate(usable_tables, start=1):
        worksheet = workbook.create_sheet(title=f"Table {index}")
        rows = table.df.fillna("").values.tolist()
        for row in rows:
            worksheet.append([str(value) for value in row])

    workbook.save(output_path)


def _convert_docx_to_pptx(input_path: Path, output_path: Path) -> None:
    try:
        from docx import Document
        from pptx import Presentation
    except ModuleNotFoundError as exc:
        raise ConversionError("missing_tool") from exc

    output_path.unlink(missing_ok=True)
    document = Document(str(input_path))
    paragraphs = [
        paragraph
        for paragraph in document.paragraphs
        if paragraph.text.strip()
    ]
    if not paragraphs:
        raise ConversionError("empty_output")

    presentation = Presentation()
    has_heading_1 = any(
        (paragraph.style.name if paragraph.style else "") == "Heading 1"
        for paragraph in paragraphs
    )
    if has_heading_1:
        _build_heading_slides(presentation, paragraphs)
    else:
        _build_chunked_slides(presentation, paragraphs)

    presentation.save(output_path)


def _add_slide(presentation, title: str, bullet_texts: list[str]) -> None:
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = title or "Document"
    placeholder = slide.placeholders[1]
    text_frame = placeholder.text_frame
    text_frame.clear()
    for index, text in enumerate(bullet_texts or [""]):
        paragraph = text_frame.paragraphs[0] if index == 0 else text_frame.add_paragraph()
        paragraph.text = text
        paragraph.level = 0


def _build_heading_slides(presentation, paragraphs: list) -> None:
    current_title: str | None = None
    current_body: list[str] = []

    for paragraph in paragraphs:
        style_name = paragraph.style.name if paragraph.style else ""
        text = paragraph.text.strip()
        if style_name == "Heading 1":
            if current_title is not None:
                _add_slide(presentation, current_title, current_body)
            current_title = text
            current_body = []
        elif current_title is None:
            current_title = "Document"
            current_body = [text]
        else:
            current_body.append(text)

    if current_title is not None:
        _add_slide(presentation, current_title, current_body)


def _build_chunked_slides(presentation, paragraphs: list, chunk_size: int = 5) -> None:
    texts = [paragraph.text.strip() for paragraph in paragraphs]
    for index in range(0, len(texts), chunk_size):
        chunk = texts[index:index + chunk_size]
        title = f"Document Part {index // chunk_size + 1}"
        _add_slide(presentation, title, chunk)
