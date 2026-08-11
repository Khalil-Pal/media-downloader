"""End-to-end Tier 2 conversion verification for the built Docker image."""
from __future__ import annotations

import asyncio
import os
import zipfile
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import AsyncMock
from urllib.parse import unquote, urlparse

import fitz
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

import services.converter as converter
from config.settings import settings
from handlers.convert_handler import handle_convertible_file
from services import db
from services.converter import (
    ConversionError,
    cleanup_conversion_session,
    convert_file,
    create_conversion_session,
    friendly_conversion_error,
)

PDF_MIME = "application/pdf"
DOCX_MIME = "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def session_listing() -> list[str]:
    if not settings.download_path.exists():
        return []
    return sorted(
        path.name
        for path in settings.download_path.iterdir()
        if path.is_dir()
    )


def make_text_pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page()
    lines = [
        "Sandy Squirrel Tier 2 Docker verification",
        "This PDF has an extractable text layer.",
        "LibreOffice should return a valid non-empty document.",
    ]
    for index, text in enumerate(lines):
        page.insert_text((72, 72 + index * 24), text, fontsize=12)
    document.save(path)
    document.close()


def make_table_pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 48), "Sandy Squirrel usage table", fontsize=14)
    rows = [
        ("Type", "Count", "Status"),
        ("Video", "10", "Ready"),
        ("Audio", "7", "Ready"),
        ("Document", "4", "Ready"),
    ]
    for row_index, row in enumerate(rows):
        y = 90 + row_index * 28
        for x, value in zip((72, 230, 380), row):
            page.insert_text((x, y), value, fontsize=12)
    document.save(path)
    document.close()


def make_image_only_pdf(path: Path) -> None:
    document = fitz.open()
    page = document.new_page()
    pixmap = fitz.Pixmap(fitz.csRGB, (0, 0, 500, 220), 0)
    pixmap.clear_with(225)
    page.insert_image(fitz.Rect(48, 48, 548, 268), pixmap=pixmap)
    document.save(path)
    document.close()


def make_heading_docx(path: Path) -> None:
    document = Document()
    document.add_heading("Download media", level=1)
    document.add_paragraph("Paste a supported media URL.")
    document.add_paragraph("Choose a quality.", style="List Bullet")
    document.add_heading("Convert files", level=1)
    document.add_paragraph("Send the source as a document.")
    document.add_paragraph("Choose a target format.", style="List Bullet")
    document.save(path)


def make_plain_docx(path: Path) -> None:
    document = Document()
    for index in range(12):
        document.add_paragraph(f"Plain paragraph {index + 1}")
    document.save(path)


def validate_docx(path: Path) -> str:
    if not zipfile.is_zipfile(path):
        raise AssertionError(f"{path.name} is not a valid OOXML archive")
    document = Document(str(path))
    return f"bytes={path.stat().st_size} paragraphs={len(document.paragraphs)}"


def validate_pptx(path: Path) -> str:
    if not zipfile.is_zipfile(path):
        raise AssertionError(f"{path.name} is not a valid OOXML archive")
    presentation = Presentation(str(path))
    if not presentation.slides:
        raise AssertionError(f"{path.name} has no slides")
    return f"bytes={path.stat().st_size} slides={len(presentation.slides)}"


def validate_markdown(path: Path) -> str:
    text = path.read_text(encoding="utf-8")
    if "Sandy Squirrel" not in text:
        raise AssertionError("Markdown output does not contain the source text")
    return f"bytes={path.stat().st_size} chars={len(text)}"


def validate_xlsx(path: Path) -> str:
    workbook = load_workbook(path, read_only=True)
    try:
        if workbook.sheetnames != ["Table 1"]:
            raise AssertionError(f"Unexpected sheets: {workbook.sheetnames}")
        values = list(workbook["Table 1"].values)
        if not values or values[0] != ("Type", "Count", "Status"):
            raise AssertionError(f"Unexpected first row: {values[:1]}")
        return (
            f"bytes={path.stat().st_size} "
            f"sheets={workbook.sheetnames} rows={len(values)}"
        )
    finally:
        workbook.close()


async def run_success_case(
    name: str,
    source_name: str,
    maker,
    target: str,
    mime_type: str,
    validator,
) -> None:
    before = session_listing()
    session = create_conversion_session()
    print(f"CASE {name}: before={before} session={session.name}")
    try:
        source = session / source_name
        maker(source)
        print(
            f"CASE {name}: command=await convert_file("
            f"'{source.name}', '{target}', '{mime_type}')"
        )
        output = await convert_file(source, target, mime_type)
        if not output.exists() or output.stat().st_size == 0:
            raise AssertionError("Conversion output is missing or empty")
        print(f"CASE {name}: exit_code=0 output={output.name} {validator(output)}")
    except Exception:
        print(f"CASE {name}: exit_code=1")
        raise
    finally:
        cleanup_conversion_session(session)
        after = session_listing()
        print(
            f"CASE {name}: after={after} "
            f"cleaned={session.name not in after}"
        )
        if before != after:
            raise AssertionError(f"{name} left an orphaned session directory")


async def run_failure_case(
    name: str,
    source_name: str,
    maker,
    target: str,
    expected_code: str,
) -> None:
    before = session_listing()
    session = create_conversion_session()
    print(f"CASE {name}: before={before} session={session.name}")
    try:
        source = session / source_name
        maker(source)
        print(
            f"CASE {name}: command=await convert_file("
            f"'{source.name}', '{target}', '{PDF_MIME}')"
        )
        try:
            await convert_file(source, target, PDF_MIME)
        except ConversionError as exc:
            print(
                f"CASE {name}: exit_code=expected_failure code={exc.code} "
                f"friendly={friendly_conversion_error(exc)}"
            )
            if exc.code != expected_code:
                raise AssertionError(
                    f"Expected {expected_code}, received {exc.code}"
                ) from exc
        else:
            raise AssertionError(f"{name} unexpectedly produced an output")
    finally:
        cleanup_conversion_session(session)
        after = session_listing()
        print(
            f"CASE {name}: after={after} "
            f"cleaned={session.name not in after}"
        )
        if before != after:
            raise AssertionError(f"{name} left an orphaned session directory")


def profile_uri_to_path(value: str) -> Path:
    parsed = urlparse(value)
    path = unquote(parsed.path)
    if os.name == "nt" and path.startswith("/") and ":" in path[:4]:
        path = path[1:]
    return Path(path)


async def run_libreoffice_concurrency_case() -> None:
    sessions = [create_conversion_session(), create_conversion_session()]
    before = [
        name for name in session_listing()
        if name not in {session.name for session in sessions}
    ]
    profile_paths: list[Path] = []
    active_processes = 0
    max_active_processes = 0
    original_run_process = converter._run_process

    async def observed_run_process(command, tool_name, timeout):
        nonlocal active_processes, max_active_processes
        profile_arg = next(
            item for item in command if item.startswith("-env:UserInstallation=")
        )
        profile_path = profile_uri_to_path(profile_arg.split("=", 1)[1])
        profile_paths.append(profile_path)
        if not profile_path.exists():
            raise AssertionError("LibreOffice profile did not exist before launch")
        active_processes += 1
        max_active_processes = max(max_active_processes, active_processes)
        print(
            f"CASE libreoffice_concurrent: process_start "
            f"profile={profile_path} active={active_processes}"
        )
        try:
            await original_run_process(command, tool_name, timeout)
        finally:
            active_processes -= 1

    converter._run_process = observed_run_process
    print(
        "CASE libreoffice_concurrent: "
        f"before={before} sessions={[session.name for session in sessions]}"
    )
    try:
        sources: list[Path] = []
        for index, session in enumerate(sessions, start=1):
            source = session / f"concurrent-{index}.pdf"
            make_text_pdf(source)
            sources.append(source)

        print(
            "CASE libreoffice_concurrent: "
            "command=await asyncio.gather(two PDF-to-DOCX convert_file calls)"
        )
        outputs = await asyncio.gather(
            *(convert_file(source, "docx", PDF_MIME) for source in sources)
        )
        details = [validate_docx(output) for output in outputs]
        if len(set(profile_paths)) != 2:
            raise AssertionError(f"Profiles were not unique: {profile_paths}")
        if max_active_processes < 2:
            raise AssertionError("LibreOffice processes did not overlap")
        if any(path.exists() for path in profile_paths):
            raise AssertionError("LibreOffice profile directory was not removed")
        print(
            "CASE libreoffice_concurrent: "
            f"exit_code=0 max_active={max_active_processes} "
            f"profiles={[str(path) for path in profile_paths]} outputs={details}"
        )
    except Exception:
        print("CASE libreoffice_concurrent: exit_code=1")
        raise
    finally:
        converter._run_process = original_run_process
        for session in sessions:
            cleanup_conversion_session(session)
        after = session_listing()
        print(
            "CASE libreoffice_concurrent: "
            f"after={after} cleaned={before == after}"
        )
        if before != after:
            raise AssertionError(
                "Concurrent conversions left orphaned session directories"
            )


async def run_mode_gate_case() -> None:
    db._pool = None
    db._memory_users.clear()
    user_id = 88001
    await db.set_user_lang(user_id, "en")
    await db.set_user_mode(user_id, "downloader")

    for file_name, mime_type in (
        ("sample.pdf", PDF_MIME),
        ("sample.docx", DOCX_MIME),
    ):
        message = SimpleNamespace(
            from_user=SimpleNamespace(id=user_id),
            document=SimpleNamespace(
                file_id="verification-file",
                file_name=file_name,
                mime_type=mime_type,
                file_size=1024,
            ),
            video=None,
            audio=None,
            answer=AsyncMock(),
        )
        await handle_convertible_file(message)
        response = message.answer.await_args.args[0]
        if "Converter mode" not in response:
            raise AssertionError(
                f"{file_name} was not gated in Downloader mode: {response}"
            )
        print(
            f"CASE mode_gate_{Path(file_name).suffix[1:]}: "
            f"exit_code=0 response={response!r}"
        )


async def main() -> None:
    if os.getenv("TIER2_VERIFY") != "1":
        raise RuntimeError("Set TIER2_VERIFY=1 to run the destructive-free Docker check")
    if settings.download_path.name != "sandy-tier2-verification":
        raise RuntimeError(
            "DOWNLOAD_PATH must end with 'sandy-tier2-verification'"
        )
    settings.download_path.mkdir(parents=True, exist_ok=True)
    if any(settings.download_path.iterdir()):
        raise RuntimeError("The verification DOWNLOAD_PATH must start empty")
    if settings.max_concurrent_conversions < 2:
        raise RuntimeError("MAX_CONCURRENT_CONVERSIONS must be at least 2")

    print(f"DOWNLOAD_PATH={settings.download_path}")
    print(f"INITIAL_SESSION_LISTING={session_listing()}")
    print(f"MAX_CONCURRENT_CONVERSIONS={settings.max_concurrent_conversions}")
    print(
        f"MAX_CONVERT_FILE_SIZE_MB={settings.max_convert_file_size_mb}"
    )

    await run_success_case(
        "pdf_to_docx",
        "text.pdf",
        make_text_pdf,
        "docx",
        PDF_MIME,
        validate_docx,
    )
    await run_success_case(
        "pdf_to_pptx",
        "text.pdf",
        make_text_pdf,
        "pptx",
        PDF_MIME,
        validate_pptx,
    )
    await run_success_case(
        "pdf_to_markdown",
        "text.pdf",
        make_text_pdf,
        "md",
        PDF_MIME,
        validate_markdown,
    )
    await run_success_case(
        "pdf_to_excel",
        "table.pdf",
        make_table_pdf,
        "xlsx",
        PDF_MIME,
        validate_xlsx,
    )
    await run_success_case(
        "docx_to_pptx_heading_1",
        "heading.docx",
        make_heading_docx,
        "pptx",
        DOCX_MIME,
        validate_pptx,
    )
    await run_success_case(
        "docx_to_pptx_no_heading",
        "plain.docx",
        make_plain_docx,
        "pptx",
        DOCX_MIME,
        validate_pptx,
    )
    await run_failure_case(
        "scanned_pdf_rejection",
        "image-only.pdf",
        make_image_only_pdf,
        "md",
        "scanned_pdf",
    )
    await run_failure_case(
        "zero_tables_rejection",
        "text-only.pdf",
        make_text_pdf,
        "xlsx",
        "no_tables",
    )
    await run_libreoffice_concurrency_case()
    await run_mode_gate_case()

    final_listing = session_listing()
    print(f"FINAL_SESSION_LISTING={final_listing}")
    if final_listing:
        raise AssertionError(f"Orphaned sessions remain: {final_listing}")
    print("TIER2_DOCKER_VERIFICATION=PASS")


if __name__ == "__main__":
    asyncio.run(main())
